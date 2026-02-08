#!/usr/bin/env python3
"""
Simple script to create PowerPoint presentations from Python code
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

class SimplePresentation:
    """Easy PowerPoint creation"""
    
    def __init__(self, title="My Presentation", output_file="output.pptx"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.output_file = output_file
        self.title = title
        
    def add_title_slide(self, title, subtitle=""):
        """Add a title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add background color (light blue)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(173, 216, 230)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(28)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    
    def add_content_slide(self, title, content_list):
        """Add a slide with title and bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, content in enumerate(content_list):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = content
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    def add_image_slide(self, title, image_path, position=(2, 1.5), width=6):
        """Add a slide with title and image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Add image
        try:
            slide.shapes.add_picture(image_path, Inches(position[0]), Inches(position[1]), width=Inches(width))
        except Exception as e:
            print(f"Error adding image: {e}")
    
    def save(self, output_dir="Slides/output"):
        """Save the presentation"""
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        filepath = Path(output_dir) / self.output_file
        self.prs.save(str(filepath))
        print(f"✅ Saved to: {filepath}")


# Example usage:
if __name__ == "__main__":
    # Create a presentation
    pres = SimplePresentation(title="My Biology Lesson", output_file="biology_lesson.pptx")
    
    # Add title slide
    pres.add_title_slide("Cell Structure", "Year 10 Biology")
    
    # Add content slides
    pres.add_content_slide(
        "What is a Cell?",
        [
            "• The basic unit of life",
            "• All living organisms are made of cells",
            "• Cells come from pre-existing cells",
            "• Can be prokaryotic or eukaryotic"
        ]
    )
    
    pres.add_content_slide(
        "Eukaryotic vs Prokaryotic",
        [
            "Eukaryotic:",
            "  - Have a nucleus",
            "  - Complex organelles",
            "  - Found in animals, plants, fungi",
            "",
            "Prokaryotic:",
            "  - No nucleus",
            "  - Simple structure",
            "  - Bacteria and archaea"
        ]
    )
    
    pres.add_content_slide(
        "Key Organelles",
        [
            "• Nucleus: contains DNA",
            "• Mitochondria: energy production",
            "• Ribosomes: protein synthesis",
            "• Endoplasmic Reticulum: transport",
            "• Golgi Apparatus: packaging"
        ]
    )
    
    # Save
    pres.save()
