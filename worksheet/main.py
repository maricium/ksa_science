#!/usr/bin/env python3
"""
Simplified worksheet generator - generates worksheets for all lessons
"""

import openai
import pandas as pd
from docx import Document
from pptx import Presentation
import re
import os
from pathlib import Path

# Configuration
openai.api_key = "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"

def extract_lesson_data(excel_file):
    """Extract lesson data from Excel file"""
    try:
        df = pd.read_excel(excel_file, skiprows=4)  # Skip header rows
        lessons = []
        
        for _, row in df.iterrows():
            if pd.notna(row.iloc[1]) and str(row.iloc[1]).strip():  # Check lesson title
                lesson = {
                    'code': str(row.iloc[0]).strip() if pd.notna(row.iloc[0]) else '',
                    'title': str(row.iloc[1]).strip(),
                    'knowledge': str(row.iloc[2]).strip() if pd.notna(row.iloc[2]) else '',
                    'skills': str(row.iloc[4]).strip() if pd.notna(row.iloc[4]) else ''
                }
                lessons.append(lesson)
        
        return lessons
    except Exception as e:
        print(f"Error reading {excel_file}: {e}")
        return []

def extract_objectives_from_pptx(lesson_code, mymastery_path="../mymastery"):
    """Extract objectives from slide 6 of matching PowerPoint file"""
    try:
        # Find matching PowerPoint file
        pptx_files = list(Path(mymastery_path).glob("*.pptx"))
        
        # Look for files that start with the lesson code
        matching_files = [f for f in pptx_files if f.stem.startswith(lesson_code)]
        
        if not matching_files:
            print(f"No PowerPoint file found for lesson {lesson_code}")
            return []
        
        # Use the first matching file
        pptx_file = matching_files[0]
        print(f"📊 Extracting objectives from: {pptx_file.name}")
        
        # Load presentation
        prs = Presentation(str(pptx_file))
        
        # Check if we have at least 6 slides (slide 6 is index 5)
        if len(prs.slides) < 6:
            print(f"Warning: {pptx_file.name} has only {len(prs.slides)} slides, expected at least 6")
            return []
        
        # Extract text from slide 6 (index 5)
        slide_6 = prs.slides[5]
        objectives = []
        
        # Extract text from all shapes in slide 6
        for shape in slide_6.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Clean up the text and split into objectives
                    lines = text.split('\n')
                    for line in lines:
                        line = line.strip()
                        if line:
                            # Check for bullet points, numbered items, or plain objectives
                            if (line.startswith('•') or line.startswith('-') or 
                                line.startswith('*') or re.match(r'^\d+\.', line)):
                                # Remove bullet points and clean up
                                clean_line = re.sub(r'^[•\-\*\d+\.\s]+', '', line).strip()
                                if clean_line:
                                    objectives.append(clean_line)
                            # Also check for plain objectives (non-empty lines that aren't headers)
                            elif (line and not line.startswith('Following this lesson') and 
                                  not line.startswith('Key Words') and 
                                  not line.startswith('Objectives') and
                                  not line.startswith('electrostatic force') and
                                  not line.startswith('lattice') and
                                  not line.startswith('ionic') and
                                  len(line) > 15 and  # Longer text is more likely to be an objective
                                  any(word in line.lower() for word in ['list', 'describe', 'explain', 'identify', 'compare', 'analyze', 'evaluate', 'calculate', 'determine'])):  # Action words typical of objectives
                                objectives.append(line)
        
        print(f"✅ Found {len(objectives)} objectives in slide 6")
        return objectives
        
    except Exception as e:
        print(f"Error extracting objectives from PowerPoint: {e}")
        return []

def generate_worksheet_with_ai(lesson, pptx_objectives=None):
    """Generate worksheet using OpenAI API"""
    
    # Build context from lesson data
    context = f"""
LESSON INFORMATION:
Title: {lesson['title']}

Knowledge students need to know:
{lesson['knowledge']}

Skills students need to demonstrate:
{lesson['skills']}
"""
    
    if pptx_objectives:
        context += f"\n\nLearning objectives from the lesson:\n"
        context += "\n".join([f"- {obj}" for obj in pptx_objectives])
    
    prompt = f"""{context}

CREATE A GCSE CHEMISTRY WORKSHEET FOR YEAR 10 STUDENTS

CRITICAL REQUIREMENTS:

1. COVER THE LESSON CONTENT COMPLETELY WITH CREATIVE QUESTIONS
   - Questions must test what students need to KNOW from the lesson
   - Questions must test what students need to be able to DO from the lesson
   - Use the knowledge and skills provided above to create questions
   - Cover ALL the learning objectives from the PowerPoint
   - Add creative elements to make questions engaging BUT stay within AQA GCSE Science content only

2. CREATIVE BUT CURRICULUM-BOUND QUESTIONS
   - Questions can be creative and engaging for Year 10 students
   - Use interesting contexts and scenarios relevant to GCSE Science
   - Include practical applications that students can relate to
   - ONLY use content that is explicitly in the AQA GCSE Chemistry specification
   - Do NOT use examples like "salt melting ice" unless specifically in the spec
   - Make questions interesting while testing the lesson content

3. STRUCTURE (12 questions total):
   
   **Section A - Knowledge Recall (Q1-4)**
   - Test the knowledge points from the lesson
   - Use fill-in-the-blank or simple recall questions
   - Cover the key facts students need to know
   - For single word gaps, use short dashes (______)
   - For sentence completion, use "..." at the end with "Finish the sentence"
   - Fill-in-the-blank questions should NOT show mark allocations
   
   **Section B - Understanding (Q5-10)**
   - Test the skills from the lesson
   - "Explain why..." questions based on lesson content
   - Cover what students need to be able to do
   
   **Section C - Extended Response (Q11-12)**
   - Question 11: Explanation question based on lesson content (3-4 marks)
   - Question 12: EITHER extended explanation question OR "explain why the student is wrong" question (4-5 marks)
   - For "student is wrong" questions: Present a common misconception, then ask students to explain why it's incorrect
   - Focus ONLY on the lesson content provided above
   - NO application questions - only explanations and understanding
   
   TOTAL MARKS: 28 maximum

3. SCAFFOLDING & SUPPORT:
   - Include sentence starters for extended answers
   - Provide key vocabulary lists where helpful
   - Add hints in italics for challenging questions
   - Use clear command words students understand

4. MAKE IT ENGAGING:
   - Use real-world contexts (cooking, technology, everyday life)
   - Include surprising or interesting facts
   - Vary question types (multiple choice, gap fill, extended response, calculations)
   - Add an optional challenge question for high achievers

5. PROFESSIONAL MARK SCHEME:
   - Detailed model answers with mark allocation
   - Show exactly what earns each mark
   - Include common misconceptions to avoid
   - Alternative acceptable answers where relevant

6. FORMATTING:
   - Clear section headings (NO mark allocations shown in section titles)
   - Appropriate answer space indicated (lines/boxes)
   - Professional layout suitable for printing
   - Individual question marks shown, but NOT section totals

DO NOT:
- Repeat the lesson title in every question
- Use vague phrases like "solve a problem" or "practical situation" without specifying what
- Create questions that could apply to any topic
- Make 1-mark questions that need multiple points
- Write questions without clear answers
- Include A-level or university-level chemistry concepts
- Use complex bonding theories or mechanisms
- Ask for detailed explanations beyond GCSE level
- Include advanced calculations or complex formulas
- Use terminology beyond GCSE specification
- Show mark allocations for fill-in-the-blank questions
- Go outside AQA GCSE Science specification content
- Include topics not covered in GCSE Chemistry
- Use examples not explicitly in the AQA specification (like "salt melting ice")
- Include application questions - focus on explanations and understanding only

EXAMPLES OF GOOD vs BAD:

❌ BAD (Too Advanced): "Explain the quantum mechanical basis of ionic bonding and how electron orbitals overlap to form ionic compounds. (4 marks)"

❌ BAD (Too Vague): "Explain how ionic bonding applies in a practical situation. (2 marks)"

❌ BAD (Too Complex for Year 10): "Compare the lattice energies of different ionic compounds and explain how this affects their properties. (4 marks)"

❌ BAD (Too Difficult): "Analyze the relationship between ionic radius and melting point in different ionic compounds. (3 marks)"

✅ GOOD (Creative + AQA GCSE): "A student is investigating ionic compounds. Explain why sodium chloride has a high melting point of 801°C. (3 marks)"

✅ GOOD (Creative + AQA GCSE): "Explain why sodium chloride conducts electricity when molten but not when solid. (4 marks)"

✅ GOOD (Fill-in-the-blank): "Complete: Ionic compounds have ______ melting points because ______ forces are ______."

✅ GOOD (Creative + AQA GCSE): "During a practical, students observe that ionic compounds are brittle. Explain why this happens. (2 marks)"

✅ GOOD (Student Wrong Question): "A student says 'Ionic compounds conduct electricity because they have free electrons.' Explain why this student is wrong. (4 marks)"

Now create the complete worksheet with mark scheme."""

    try:
        client = openai.OpenAI(api_key=openai.api_key)
        response = client.chat.completions.create(
            model="gpt-4",  # Use gpt-4 for best results, or "gpt-3.5-turbo" 
            messages=[
                {"role": "system", "content": "You are an outstanding GCSE Chemistry teacher with 20 years of experience. You create engaging, clear, and pedagogically sound worksheets that students actually enjoy completing. You understand how to scaffold learning and make abstract concepts concrete."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=3000  # Increased for detailed worksheets
        )
        
        return response.choices[0].message.content
        
    except Exception as e:
        print(f"❌ OpenAI API error: {e}")
        return None

def generate_worksheet(lesson, subject="Chemistry"):
    """Generate a worksheet using OpenAI with PowerPoint objectives"""
    print(f"Generating worksheet for: {lesson['title']}")
    print("🤖 Using OpenAI API for generation")
    
    # Extract objectives from PowerPoint
    pptx_objectives = extract_objectives_from_pptx(lesson['code'])
    
    # Try AI generation first
    ai_content = generate_worksheet_with_ai(lesson, pptx_objectives)
    
    if ai_content:
        return ai_content
    else:
        print("⚠️ AI generation failed, falling back to template")
        return create_template_worksheet(lesson, pptx_objectives)

def create_template_worksheet(lesson, pptx_objectives=None):
    """Create a worksheet based on specific knowledge and skills from the lesson and PowerPoint objectives"""
    knowledge_points = lesson['knowledge'].split('\n') if lesson['knowledge'] else []
    skills_points = lesson['skills'].split('\n') if lesson['skills'] else []
    
    # Filter out empty lines and clean up
    knowledge_points = [point.strip() for point in knowledge_points if point.strip()]
    skills_points = [point.strip() for point in skills_points if point.strip()]
    
    # Create questions based on actual knowledge, skills, and PowerPoint objectives
    knowledge_questions = []
    skills_questions = []
    analysis_questions = []  # Initialize analysis questions list
    objective_question_mapping = []  # Track which questions map to which objectives
    
    # Use PowerPoint objectives if available, otherwise fall back to lesson content
    if pptx_objectives:
        print(f"🎯 Using {len(pptx_objectives)} objectives from PowerPoint slide 6")
        
        # Create unique questions for each objective
        used_objectives = set()
        
        # Generate knowledge questions (Section A) - use first 4 objectives or create unique ones
        for i in range(4):
            question_num = i + 1
            if i < len(pptx_objectives) and pptx_objectives[i] not in used_objectives:
                objective = pptx_objectives[i]
                knowledge_questions.append(f"{question_num}. {objective} (1 mark)")
                objective_question_mapping.append(f"• {objective} → met Q{question_num}")
                used_objectives.add(objective)
            else:
                # Create unique fallback questions
                unique_questions = [
                    f"State the main properties of {lesson['title'].lower()}.",
                    f"Define a key characteristic of {lesson['title'].lower()}.",
                    f"Identify the primary feature of {lesson['title'].lower()}.",
                    f"Recall an important aspect of {lesson['title'].lower()}."
                ]
                if question_num <= len(unique_questions):
                    knowledge_questions.append(f"{question_num}. {unique_questions[question_num-1]} (1 mark)")
                    objective_question_mapping.append(f"• {unique_questions[question_num-1]} → met Q{question_num}")
        
        # Generate skills questions (Section B) - use next objectives or create unique ones
        for i in range(4):
            question_num = i + 5
            objective_index = i + 4  # Start from 4th objective
            
            if (objective_index < len(pptx_objectives) and 
                pptx_objectives[objective_index] not in used_objectives):
                objective = pptx_objectives[objective_index]
                skills_questions.append(f"{question_num}. {objective} (2 marks)")
                objective_question_mapping.append(f"• {objective} → met Q{question_num}")
                used_objectives.add(objective)
            else:
                # Create unique fallback questions
                unique_skills = [
                    f"Apply your understanding of {lesson['title'].lower()} to solve a problem.",
                    f"Calculate using principles from {lesson['title'].lower()}.",
                    f"Explain how {lesson['title'].lower()} applies in a practical situation.",
                    f"Compare different aspects of {lesson['title'].lower()}."
                ]
                if (question_num - 5) < len(unique_skills):
                    skills_questions.append(f"{question_num}. {unique_skills[question_num-5]} (2 marks)")
                    objective_question_mapping.append(f"• {unique_skills[question_num-5]} → met Q{question_num}")
        
        # Generate analysis questions (Section C) - use remaining objectives or create unique ones
        for i in range(4):
            question_num = i + 9
            objective_index = i + 8  # Start from 8th objective
            
            if (objective_index < len(pptx_objectives) and 
                pptx_objectives[objective_index] not in used_objectives):
                objective = pptx_objectives[objective_index]
                analysis_questions.append(f"{question_num}. {objective} (3-4 marks)")
                objective_question_mapping.append(f"• {objective} → met Q{question_num}")
                used_objectives.add(objective)
            else:
                # Create unique fallback questions
                unique_analysis = [
                    f"Analyze the relationship between different concepts in {lesson['title'].lower()}.",
                    f"Evaluate the importance of {lesson['title'].lower()} in chemistry.",
                    f"Compare and contrast {lesson['title'].lower()} with related topics.",
                    f"Synthesize your knowledge of {lesson['title'].lower()} to explain a complex scenario."
                ]
                if (question_num - 9) < len(unique_analysis):
                    analysis_questions.append(f"{question_num}. {unique_analysis[question_num-9]} (3-4 marks)")
                    objective_question_mapping.append(f"• {unique_analysis[question_num-9]} → met Q{question_num}")
    
    else:
        # Fall back to original lesson content with unique questions
        for i, knowledge in enumerate(knowledge_points[:4]):
            if knowledge:
                question_num = i + 1
                knowledge_questions.append(f"{question_num}. {knowledge} (1 mark)")
                objective_question_mapping.append(f"• {knowledge} → met Q{question_num}")
        
        for i, skill in enumerate(skills_points[:4]):
            if skill:
                question_num = i + 5
                skills_questions.append(f"{question_num}. {skill} (2 marks)")
                objective_question_mapping.append(f"• {skill} → met Q{question_num}")
        
        # Fill remaining questions if needed with unique content
        while len(knowledge_questions) < 4:
            question_num = len(knowledge_questions) + 1
            unique_questions = [
                f"State the main properties of {lesson['title'].lower()}.",
                f"Define a key characteristic of {lesson['title'].lower()}.",
                f"Identify the primary feature of {lesson['title'].lower()}.",
                f"Recall an important aspect of {lesson['title'].lower()}."
            ]
            if question_num <= len(unique_questions):
                knowledge_questions.append(f"{question_num}. {unique_questions[question_num-1]} (1 mark)")
                objective_question_mapping.append(f"• {unique_questions[question_num-1]} → met Q{question_num}")
        
        while len(skills_questions) < 4:
            question_num = len(skills_questions) + 5
            unique_skills = [
                f"Apply your understanding of {lesson['title'].lower()} to solve a problem.",
                f"Calculate using principles from {lesson['title'].lower()}.",
                f"Explain how {lesson['title'].lower()} applies in a practical situation.",
                f"Compare different aspects of {lesson['title'].lower()}."
            ]
            if (question_num - 5) < len(unique_skills):
                skills_questions.append(f"{question_num}. {unique_skills[question_num-5]} (2 marks)")
                objective_question_mapping.append(f"• {unique_skills[question_num-5]} → met Q{question_num}")
        
        # Create analysis questions (Section C)
        for i in range(4):
            question_num = i + 9
            unique_analysis = [
                f"Analyze the relationship between different concepts in {lesson['title'].lower()}.",
                f"Evaluate the importance of {lesson['title'].lower()} in chemistry.",
                f"Compare and contrast {lesson['title'].lower()} with related topics.",
                f"Synthesize your knowledge of {lesson['title'].lower()} to explain a complex scenario."
            ]
            analysis_questions.append(f"{question_num}. {unique_analysis[i]} (3-4 marks)")
            objective_question_mapping.append(f"• {unique_analysis[i]} → met Q{question_num}")
    
    # Combine all questions
    all_questions = knowledge_questions + skills_questions + analysis_questions
    
    # Create simple answer-based mark scheme for student self-correction
    mark_scheme = []
    for i, question in enumerate(all_questions, 1):
        question_text = question.split('.', 1)[1].split('(')[0].strip()
        marks = question.split('(')[-1].replace(')', '').strip()
        
        # Generate simple answers based on question content and lesson topic
        topic_lower = lesson['title'].lower()
        
        if i <= 4:  # Knowledge questions (1 mark) - simple factual answers
            if "list" in question_text.lower() and "properties" in question_text.lower():
                answer = "High melting/boiling point, conducts electricity when molten/dissolved, brittle, soluble in water"
            elif "describe" in question_text.lower() and "structure" in question_text.lower():
                answer = "Regular lattice of positive and negative ions held together by strong electrostatic forces"
            elif "explain" in question_text.lower() and "reference" in question_text.lower():
                answer = "Strong electrostatic forces between ions require lots of energy to break, explaining high melting point"
            else:
                answer = f"Key concept from {topic_lower} - students should state main properties or structure"
            
            mark_scheme.append(f"{i}. {answer}")
            
        elif i <= 8:  # Skills questions (2 marks) - application examples
            if "apply" in question_text.lower() and "problem" in question_text.lower():
                answer = "Use knowledge of ionic bonding to predict properties - e.g. high melting point means strong forces"
            elif "calculate" in question_text.lower():
                answer = "Apply formula or use ionic bonding principles to solve numerical problem"
            elif "explain" in question_text.lower() and "practical" in question_text.lower():
                answer = "Real-world example - e.g. salt dissolving in water due to ion separation"
            elif "compare" in question_text.lower():
                answer = "Contrast ionic substances with covalent/metallic - different bonding = different properties"
            else:
                answer = f"Apply {topic_lower} knowledge to solve problem or explain scenario"
            
            mark_scheme.append(f"{i}. {answer}")
            
        else:  # Analysis questions (3-4 marks) - deeper thinking
            if "analyze" in question_text.lower() and "relationship" in question_text.lower():
                answer = "Link structure to properties - regular lattice → high melting point, free ions → conductivity"
            elif "evaluate" in question_text.lower() and "importance" in question_text.lower():
                answer = "Essential for life (salts), industry (electrolysis), daily life (table salt)"
            elif "compare and contrast" in question_text.lower():
                answer = "Ionic vs covalent vs metallic bonding - different structures lead to different properties"
            elif "synthesize" in question_text.lower():
                answer = "Combine multiple concepts to explain complex situation - structure + bonding + properties"
            else:
                answer = f"Deep analysis of {topic_lower} - connect multiple ideas and explain relationships"
            
            mark_scheme.append(f"{i}. {answer}")
        
        mark_scheme.append("")
    
    # Create objectives section
    objectives_section = ""
    if pptx_objectives:
        objectives_section = f"""
LEARNING OBJECTIVES FROM POWERPOINT (Slide 6):
{chr(10).join([f"• {obj}" for obj in pptx_objectives])}
"""
    
    # Create objective-to-question mapping section
    mapping_section = ""
    if objective_question_mapping:
        mapping_section = f"""

OBJECTIVE-QUESTION MAPPING:
{chr(10).join(objective_question_mapping)}
"""
    
    template = f"""GCSE Chemistry Worksheet: {lesson['title']}

Section A - Knowledge Recall (Questions 1-4)

{chr(10).join(knowledge_questions)}

Section B - Skills Application (Questions 5-8)

{chr(10).join(skills_questions)}

Section C - Analysis and Evaluation (Questions 9-12)

{chr(10).join(analysis_questions)}

---

MARK SCHEME

{chr(10).join(mark_scheme)}

---

LEARNING OBJECTIVES REFERENCED:

What students need to KNOW:
{lesson['knowledge']}

What students need to be able to DO:
{lesson['skills']}
{objectives_section}{mapping_section}
"""

    return template

def create_word_document(content, lesson):
    """Create Word document from AI content"""
    doc = Document()
    
    # Set narrow margins
    sections = doc.sections
    for section in sections:
        section.left_margin = section.right_margin = section.top_margin = section.bottom_margin = 457200  # 0.5 inch margins
    
    # Add title
    doc.add_heading(lesson['title'], 0)
    doc.add_paragraph("")
    
    # Process content
    lines = content.split('\n')
    in_questions = True
    import re
    question_pattern = re.compile(r'^\d+\.')

    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Check for mark scheme section
        if "mark scheme" in line.lower() or "marking" in line.lower():
            in_questions = False
            doc.add_page_break()
            doc.add_heading("MARK SCHEME", 1)
            continue
        
        if in_questions and question_pattern.match(line):
            # Add question
            doc.add_paragraph(line)
            
            # Check if this is a fill-in-the-blank question using regex for multiple underscores
            is_fill_in_blank = ("complete" in line.lower() or "fill" in line.lower() or 
                              re.search(r'_{3,}', line) or "gap" in line.lower())
            
            if not is_fill_in_blank:
                # Add answer lines for non-fill-in-the-blank questions
                question_num = int(line.split('.')[0])
                num_lines = 2 if question_num <= 4 else 3 if question_num <= 8 else 4
                
                for _ in range(num_lines):
                    p = doc.add_paragraph()
                    p.add_run("_" * 120)
            
            doc.add_paragraph("")
        else:
            doc.add_paragraph(line)

    return doc

def main():
    """Generate single worksheet for C4.1.3"""
    print("🔄 Generating C4.1.3 worksheet...")
    
    # Extract lessons from C4.1
    excel_file = "../Links/C4.1 All Resources/c4.1 unit overview xslx/C4.1 Unit Overview (4).xlsx"
    
    if os.path.exists(excel_file):
        lessons = extract_lesson_data(excel_file)
        print(f"📊 Found {len(lessons)} lessons in C4.1")
        
        # Find C4.1.3 lesson
        target_lesson = None
        for lesson in lessons:
            if 'C4.1.3' in lesson['code']:
                target_lesson = lesson
                break
        
        if target_lesson:
            print(f"Creating: {target_lesson['title']}")
            
            # Generate worksheet content
            content = generate_worksheet(target_lesson)
            if content:
                # Create Word document
                doc = create_word_document(content, target_lesson)
                
                # Save worksheet
                os.makedirs("worksheets/C4.1", exist_ok=True)
                filename = f"{target_lesson['code']}_{target_lesson['title'].replace(' ', '_')}_worksheet.docx"
                doc.save(f"worksheets/C4.1/{filename}")
                print(f"✅ Created: C4.1/{filename}")
            else:
                print("❌ Failed to generate worksheet content")
        else:
            print("❌ C4.1.3 lesson not found")
            print("Available C4.1 lessons:")
            for lesson in lessons:
                print(f"  {lesson['code']} - {lesson['title']}")
    else:
        print(f"❌ Excel file not found: {excel_file}")

if __name__ == "__main__":
    main()