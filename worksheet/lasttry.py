#!/usr/bin/env python3
"""Worksheet Generator - validates existing worksheets against Excel 'know' column"""

import sys
from pathlib import Path
from docx import Document
from docx.shared import Inches
import anthropic
import os
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).parent.parent))
from read_excel import get_lesson_data

def extract_slide6_objectives(lesson_code, mymastery_dir="Lesson Resources", slide_num=6):
    """Extract text from specified slide of PowerPoint file"""
    try:
        mymastery_path = Path(mymastery_dir)
        if not mymastery_path.exists():
            mymastery_path = Path("..") / mymastery_dir
        if not mymastery_path.exists():
            mymastery_path = Path("../..") / mymastery_dir
        
        if mymastery_path.is_dir():
            pptx_files = list(mymastery_path.rglob(f"{lesson_code}*.pptx"))
            pptx_files = [f for f in pptx_files if not f.name.startswith('~')]
            if not pptx_files:
                return ""
            pptx_path = pptx_files[0]
        else:
            pptx_path = mymastery_path
        
        prs = Presentation(str(pptx_path))
        slide_index = slide_num - 1
        if len(prs.slides) < slide_num:
            return ""
        
        slide_text = []
        for shape in prs.slides[slide_index].shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        return "\n".join(slide_text)
    except Exception as e:
        print(f"   ⚠️  Error extracting objectives: {e}")
        return ""

env_file = Path(__file__).parent.parent / '.env'
if env_file.exists():
    for line in env_file.read_text().splitlines():
        if line.startswith('ANTHROPIC_API_KEY='):
            os.environ['ANTHROPIC_API_KEY'] = line.split('=', 1)[1].strip('"').strip("'")

ANTHROPIC_API_KEY = os.getenv('ANTHROPIC_API_KEY') or "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"
if not ANTHROPIC_API_KEY:
    print("❌ Error: ANTHROPIC_API_KEY not found")
    exit(1)

ANSWER_LINE_LENGTH, MARGIN_SIZE, FONT_NAME = 120, 0.5, 'Century Gothic'

# Configuration for foundation/higher tier separation
STUDENT_YEAR = 10  # Change this to match your class (e.g., 7, 8, 9, 10, 11)
STUDENT_ATTAINMENT = "LPA"  # "LPA" or "HPA"
FOUNDATION_ONLY = STUDENT_ATTAINMENT.upper() == "LPA" and STUDENT_YEAR >= 10  # Foundation only for LPA Year 10+

def find_path(name):
    """Find path in current or parent directories"""
    for p in [Path(name), Path(f"../{name}"), Path(f"../../{name}")]:
        if p.exists():
            return p
    return Path(name)

def filter_foundation_content(know_content, lesson_title):
    """Filter content to foundation tier only - remove higher tier (HT) content"""
    if not FOUNDATION_ONLY:
        return know_content
    
    # Check if lesson title indicates higher tier
    if "(HT)" in lesson_title or "HT only" in lesson_title or "Higher tier" in lesson_title:
        return ""  # Skip HT-only lessons entirely
    
    # Filter out HT-specific bullet points from know content
    lines = know_content.split('\n')
    foundation_lines = []
    
    for line in lines:
        line_lower = line.lower()
        # Skip lines that explicitly mention higher tier
        if any(marker in line_lower for marker in ['(ht)', 'higher tier', 'ht only', 'higher tier only']):
            continue
        # Keep foundation content
        foundation_lines.append(line)
    
    filtered = '\n'.join(foundation_lines).strip()
    
    if not filtered or len(filtered) < 20:
        print(f"   ⚠️  Warning: After filtering for foundation tier, content is very limited or empty")
    
    return filtered

def generate_worksheet_content(lesson, pptx_objectives=None):
    """Generate worksheet content based on Excel 'know' column and PowerPoint objectives"""
    # Filter for foundation tier if needed
    know_content = lesson['know']
    if FOUNDATION_ONLY:
        print(f"   📚 Filtering for Foundation Tier only (LPA Year {STUDENT_YEAR}+)")
        know_content = filter_foundation_content(lesson['know'], lesson['title'])
        if not know_content or len(know_content) < 20:
            print(f"   ⚠️  Warning: No foundation content found for {lesson['code']}")
            return None
    
    excel_knowledge = f"LESSON: {lesson['title']}\n\nWHAT STUDENTS NEED TO KNOW (from Excel):\n{know_content}"
    if pptx_objectives:
        excel_knowledge += f"\n\nLEARNING OBJECTIVES (from PowerPoint slide 6):\n{pptx_objectives}"
    
    tier_note = "FOUNDATION TIER ONLY" if FOUNDATION_ONLY else "FOUNDATION AND HIGHER TIER"
    prompt = f"""Create a GCSE Chemistry worksheet for {tier_note}. ALL questions MUST test knowledge from:\n\n{excel_knowledge}\n\n{"IMPORTANT: This worksheet is for FOUNDATION TIER students only. Do NOT include any higher tier (HT) content. Only use content that is accessible to foundation tier students." if FOUNDATION_ONLY else ""}\n\nStructure: Section 1 (brief intro), then Section A (Fill blanks, Matching, Labelling), Section B (8-10 questions), Section C (2 extended), Mark Scheme.\n\nSection A: Part 1 (Word Bank with 10-12 keywords, then 7 fill-in-blank sentences on ONE line), Part 2 (MATCHING_TABLE_START/END with 4 term|description pairs), Part 3 (Labelling diagram with 3-4 labels). All content MUST come from "WHAT STUDENTS NEED TO KNOW" only."""
    
    try:
        print("🤖 Generating with Claude...")
        client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
        message = client.messages.create(model="claude-sonnet-4-20250514", max_tokens=8000, messages=[{"role": "user", "content": prompt}])
        return message.content[0].text
    except Exception as e:
        print(f"❌ API error: {e}")
        return None

def set_font(runs, bold=False):
    """Set Comic Sans font"""
    for run in runs:
        run.font.name = FONT_NAME
        if bold:
            run.font.bold = True

def create_word_document(lesson, ai_content):
    """Create Word document"""
    if not ai_content or not ai_content.strip():
        print("   ⚠️  Warning: AI content is empty!")
        return Document()
    
    doc = Document()
    for section in doc.sections:
        margin = Inches(MARGIN_SIZE)
        section.left_margin = section.right_margin = section.top_margin = section.bottom_margin = margin
    
    current_section, in_markscheme, in_word_bank, in_matching_table = None, False, False, False
    matching_rows, found_section = [], False
    
    lines = ai_content.split('\n')
    print(f"   📝 Processing {len(lines)} lines from AI content...")
    
    # Try to find section start in first 30 lines
    for i, line in enumerate(lines[:30]):
        line_stripped = line.strip()
        if (('Section 1' in line_stripped or 'Section A' in line_stripped) and 
            ('–' in line_stripped or '-' in line_stripped or ':' in line_stripped)):
            found_section = True
            print(f"   ✓ Found section start at line {i+1}: {line_stripped[:50]}")
            break
    
    # If no section found, start anyway (content might be formatted differently)
    if not found_section:
        print(f"   ⚠️  No section header found, starting anyway...")
        found_section = True
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Skip content before sections only if we haven't found one yet
        if not found_section:
            if (('Section 1' in line or 'Section A' in line) and 
                ('–' in line or '-' in line or ':' in line)):
                found_section = True
                print(f"   ✓ Found section start: {line[:50]}")
            else:
                continue
        
        if 'MARKSCHEME' in line.upper() or '---MARK' in line:
            doc.add_page_break()
            set_font(doc.add_heading("MARK SCHEME", 0).runs, True)
            in_markscheme = True
            continue
        
        if 'MATCHING_TABLE_START' in line:
            in_matching_table, matching_rows = True, []
            continue
        if in_matching_table:
            if 'MATCHING_TABLE_END' in line:
                table = doc.add_table(rows=len(matching_rows), cols=3)
                table.style = 'Table Grid'
                table.columns[0].width = Inches(2.5)
                table.columns[1].width = Inches(1.5)
                table.columns[2].width = Inches(3.0)
                for idx, row_data in enumerate(matching_rows):
                    parts = row_data.split('|')
                    if len(parts) == 2:
                        table.rows[idx].cells[0].text = parts[0].strip()
                        set_font(table.rows[idx].cells[0].paragraphs[0].runs, True)
                        table.rows[idx].cells[2].text = parts[1].strip()
                        set_font(table.rows[idx].cells[2].paragraphs[0].runs)
                doc.add_paragraph()
                in_matching_table = False
            else:
                matching_rows.append(line)
            continue
        
        # Section headers - check for various dash types
        if 'Section 1' in line and ('–' in line or '-' in line or ':' in line):
            set_font(doc.add_heading(line, 1).runs, True)
            current_section = '1'
            continue
        if 'Section A' in line and ('–' in line or '-' in line or ':' in line):
            set_font(doc.add_heading(line, 1).runs, True)
            current_section = 'A'
            continue
        if 'Section B' in line and ('–' in line or '-' in line or ':' in line):
            doc.add_page_break()
            set_font(doc.add_heading(line, 1).runs, True)
            current_section = 'B'
            continue
        if 'Section C' in line and ('–' in line or '-' in line or ':' in line):
            doc.add_page_break()
            set_font(doc.add_heading(line, 1).runs, True)
            current_section = 'C'
            continue
        
        if line.startswith('Part '):
            set_font(doc.add_heading(line, 2).runs, True)
            continue
        
        if 'Word Bank' in line and ':' in line:
            r = doc.add_paragraph().add_run(line)
            set_font([r], True)
            in_word_bank = True
            continue
        if in_word_bank and current_section == 'A':
            table = doc.add_table(rows=1, cols=1)
            table.style = 'Table Grid'
            set_font(table.rows[0].cells[0].paragraphs[0].runs)
            table.rows[0].cells[0].text = line
            doc.add_paragraph()
            in_word_bank = False
            continue
        
        if line and line[0].isdigit() and '. ' in line:
            parts = line.split('. ', 1)
            if len(parts) == 2:
                r = doc.add_paragraph().add_run(f"{parts[0]}. {parts[1]}")
                set_font([r], True)
                if not in_markscheme and current_section in ['B', 'C'] and '_____' not in line:
                    for _ in range(3 if current_section == 'B' else 6):
                        set_font(doc.add_paragraph("_" * ANSWER_LINE_LENGTH).runs)
                    doc.add_paragraph()
            continue
        
        if in_markscheme and ('Section' in line and ':' in line or line.startswith('Part ')):
            set_font(doc.add_heading(line, 2 if 'Section' in line else 3).runs, True)
            continue
        
        if not line.startswith('---'):
            set_font(doc.add_paragraph(line).runs)
    
    return doc

def main():
    lesson = get_lesson_data("C3.1.3")
    if not lesson:
        print("❌ Could not load lesson")
        return
    
    print(f"📋 Lesson: {lesson['title']}")
    print("📄 Extracting PowerPoint objectives...")
    lesson_resources = find_path("Lesson Resources")
    pptx_objectives = extract_slide6_objectives(lesson['code'], str(lesson_resources), 6)
    if pptx_objectives:
        print(f"   ✓ Found ({len(pptx_objectives)} chars)")
    else:
        print("   ⚠️  No objectives found")
    
    print("🤖 Generating worksheet with Anthropic...")
    content = generate_worksheet_content(lesson, pptx_objectives)
    if not content:
        print("❌ Failed to generate")
        return
    
    print("✅ Content generated")
    print(f"   📊 Content length: {len(content)} characters")
    print(f"   📄 First 200 chars: {content[:200]}...")
    
    doc = create_word_document(lesson, content)
    
    # Check if document has content
    if len(doc.paragraphs) == 0:
        print("   ⚠️  Warning: Document has no paragraphs!")
    else:
        print(f"   ✓ Document has {len(doc.paragraphs)} paragraphs")
    
    Path("worksheets").mkdir(exist_ok=True)
    # Add tier suffix to filename
    tier_suffix = "_Foundation" if FOUNDATION_ONLY else "_Higher"
    filepath = Path("worksheets") / f"{lesson['code']}_worksheet{tier_suffix}.docx"
    doc.save(str(filepath))
    print(f"\n✅ Created: {filepath}")
    print(f"   Tier: {'Foundation' if FOUNDATION_ONLY else 'Foundation & Higher'}")
    print(f"   Class: {STUDENT_ATTAINMENT} Year {STUDENT_YEAR}\n🎉 Done!")

if __name__ == "__main__":
    main()
