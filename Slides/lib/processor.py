"""Process PowerPoint templates to create lesson slides"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from datetime import datetime
from pptx.enum.text import MSO_AUTO_SIZE
from .utils import find_path
from .reader import read_markscheme
import config


def _find_original_pptx(lesson_code):
    """Find original PowerPoint file for lesson code"""
    lesson_resources = find_path(config.LESSON_RESOURCES_DIR)
    pptx_files = [f for f in lesson_resources.rglob(f"{lesson_code}*.pptx") if not f.name.startswith('~')]
    return str(pptx_files[0]) if pptx_files else None


def _extract_objectives(lesson_code):
    """Extract learning objectives from original PowerPoint"""
    pptx_path = _find_original_pptx(lesson_code)
    if not pptx_path:
        return ""
    
    try:
        prs = Presentation(pptx_path)
        # Extract objectives from slides 4-7 of source PowerPoint
        for slide_num in [4, 5, 6, 7]:
            if len(prs.slides) < slide_num:
                continue
            slide = prs.slides[slide_num - 1]
            shape_texts = [(s.top.inches, s.left.inches, [p.text.strip() for p in s.text_frame.paragraphs if p.text.strip() and all(x not in p.text.lower() for x in ['fix-it', 'customise', 'reteach', 'redrafting'])]) for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame]
            shape_texts = [st for st in shape_texts if st[2]]
            if shape_texts:
                shape_texts.sort(key=lambda x: (x[0], x[1]))
                lines = [line for _, _, paras in shape_texts for line in paras]
                filtered = []
                found_header = False
                skip = False
                
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    
                    # Check if this line contains the header
                    if 'following this lesson' in line.lower() and 'students will be able to' in line.lower():
                        found_header = True
                        filtered.append(line)
                        skip = False
                        continue
                    
                    # If we found the header, include all bullet points after it
                    if found_header:
                        if 'key words' in line.lower() or 'key word' in line.lower():
                            break  # Stop at key words section
                        # Include bullet points and objectives
                        filtered.append(line)
                    else:
                        # Before finding header, skip key words section
                        if 'key words' in line.lower() or 'key word' in line.lower():
                            skip = True
                            continue
                        if skip and len(line.split()) <= 2:
                            continue
                        if skip and len(line) > 20:
                            skip = False
                        # Only include if it looks like an objective
                        if not skip and ('following' not in line.lower() or 'students will' in line.lower()):
                            filtered.append(line)
                
                if filtered:
                    print(f"   ✓ Extracted objectives from slide {slide_num}")
                    return "\n".join(filtered)
    except Exception as e:
        print(f"⚠️  Error reading PowerPoint: {e}")
    return ""


def _remove_placeholders(slide, target_left, target_top, tolerance=0.5):
    """Remove interfering placeholders near target position"""
    for shape in list(slide.shapes):
        if hasattr(shape, 'text_frame') and shape.text_frame and shape.is_placeholder:
            if abs(shape.left.inches - target_left.inches) < tolerance and abs(shape.top.inches - target_top.inches) < tolerance:
                shape._element.getparent().remove(shape._element)


def _create_textbox(slide, pos, content, font_size, color=None, bold=False, add_border=False):
    """Create and configure textbox with content"""
    textbox = slide.shapes.add_textbox(Inches(pos['left']), Inches(pos['top']), Inches(pos['width']), Inches(pos['height']))
    textbox.text_frame.word_wrap = True
    textbox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    for margin in ['margin_left', 'margin_right', 'margin_top', 'margin_bottom']:
        setattr(textbox.text_frame, margin, Inches(0.1))
    textbox.text_frame.clear()
    
    # Add border if requested
    if add_border:
        line = textbox.line
        line.color.rgb = RGBColor(0, 0, 0)  # Black border
        line.width = Pt(1.5)  # Border width
    
    for line in content.split('\n'):
        if line.strip():
            para = textbox.text_frame.add_paragraph()
            run = para.add_run()
            run.text = line.strip()
            run.font.name = 'Century Gothic'
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color
            if bold:
                run.font.bold = True
        else:
            textbox.text_frame.add_paragraph()
    return textbox


def create_slide_from_template(template_path, lesson_data, output_path):
    """Create presentation from template with lesson data"""
    try:
        prs = Presentation(template_path)
        print(f"Loaded template with {len(prs.slides)} slides")
        
        lesson_code = lesson_data.get('lesson_code', '')
        markscheme = read_markscheme(lesson_code)
        
        # Use objectives from Excel data, fallback to extracting from PowerPoint if not available
        knowledge_obj = lesson_data.get('knowledge_objectives', '').strip()
        skill_obj = lesson_data.get('skill_objectives', '').strip()
        
        if knowledge_obj and knowledge_obj != 'nan':
            # Combine knowledge and skill objectives
            objectives_parts = []
            if knowledge_obj:
                objectives_parts.append("Knowledge Objectives:")
                objectives_parts.append(knowledge_obj)
            if skill_obj and skill_obj != 'nan':
                objectives_parts.append("\nSkill Objectives:")
                objectives_parts.append(skill_obj)
            objectives = '\n'.join(objectives_parts)
        else:
            # Fallback to extracting from original PowerPoint
            objectives = _extract_objectives(lesson_code) or f"Learning Objectives\n\nObjectives not found for {lesson_code}"
        
        exit_ticket = lesson_data.get('exit_ticket', '').strip()
        
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            print(f"   Processing slide {slide_num}...")
            
            if slide_num == config.MARKSCHEME_SLIDE and markscheme:
                pos = config.MARKSCHEME_POSITION
                _remove_placeholders(slide, Inches(pos['left']), Inches(pos['top']))
                _create_textbox(slide, pos, markscheme, config.MARKSCHEME_FONT_SIZE, RGBColor(0, 102, 0), bold=True, add_border=False)
                print(f"      ✓ Updated markscheme")
            
            if slide_num == config.OBJECTIVES_SLIDE and objectives:
                pos = config.OBJECTIVES_POSITION
                _remove_placeholders(slide, Inches(pos['left']), Inches(pos['top']))
                # Format objectives - preserve existing bullets, add to lines without them
                formatted_lines = []
                for line in objectives.split('\n'):
                    line = line.strip()
                    if not line:
                        continue
                    # If line already has a bullet or is a header, keep as is
                    if line.startswith('•') or line.startswith('-') or line.endswith(':'):
                        formatted_lines.append(line)
                    elif 'objectives' in line.lower() and line.endswith(':'):
                        formatted_lines.append(line)
                    else:
                        # Add bullet if it's a content line
                        formatted_lines.append(f'• {line}')
                formatted = '\n'.join(formatted_lines)
                textbox = _create_textbox(slide, pos, formatted, config.OBJECTIVES_FONT_SIZE)
                for para in textbox.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0)
                print(f"      ✓ Updated objectives")
            
            if slide_num == config.EXIT_TICKET_SLIDE and exit_ticket:
                pos = config.EXIT_TICKET_POSITION
                _remove_placeholders(slide, Inches(pos['left']), Inches(pos['top']))
                _create_textbox(slide, pos, exit_ticket, config.EXIT_TICKET_FONT_SIZE)
                print(f"      ✓ Updated exit ticket")
            
            if slide_num == 1:
                # Only update the lesson title, nothing else on slide 1
                for shape in slide.shapes:
                    if not hasattr(shape, 'text_frame') or not shape.text_frame:
                        continue
                    text = shape.text_frame.text.strip().lower()
                    # Only update if it's clearly a title placeholder or "Lesson Title" placeholder
                    if text in ['lesson title', 'title', ''] or (shape.is_placeholder and len(text) < 20):
                        shape.text_frame.text = lesson_data.get('lesson_title', 'Lesson Title')
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.color.rgb = RGBColor(0, 0, 0)
                        print(f"      ✓ Updated title only")
        
        prs.save(output_path)
        print(f"Created: {output_path}")
        return output_path
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        return None
