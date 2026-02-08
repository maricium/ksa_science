import os
from pathlib import Path

import requests
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation

# Load environment variables from .env file
load_dotenv()

try:
    from googleapiclient.discovery import build
except ImportError:
    build = None  # Optional dependency

# Initialize OpenAI client
try:
    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
except:
    client = None
    print("⚠️  OpenAI API key not found")


def download_chemistry_image(lesson_title):
    """Download image from Google or use placeholder."""
    if build is None:
        return "placeholder.jpg"
    try:
        service = build("customsearch", "v1", developerKey=os.getenv("GOOGLE_API_KEY"))
        res = (
            service.cse()
            .list(q=f"savemyexams {lesson_title}", cx=os.getenv("GOOGLE_CX"), searchType="image", num=1)
            .execute()
        )
        img_path = f"temp_image_{lesson_title.replace(' ', '_')}.jpg"
        with open(img_path, "wb") as f:
            f.write(requests.get(res["items"][0]["link"]).content)
        return img_path
    except:
        return "placeholder.jpg"


def extract_slide6_objectives(lesson_code, mymastery_dir="Lesson Resources", slide_num=None):
    """Extract text from objectives slide of PowerPoint file. Searches slides 4-7 if slide_num not specified."""
    try:
        mymastery_path = Path(mymastery_dir)

        # Resolve path - check current dir and parent dir
        if not mymastery_path.exists():
            mymastery_path = Path("..") / mymastery_dir
        if not mymastery_path.exists():
            mymastery_path = Path("../..") / mymastery_dir

        # If directory provided, search recursively for PowerPoint file matching lesson code
        if mymastery_path.is_dir():
            # Search recursively in subfolders (e.g., Lesson Resources/C3.1.3 Electronic Configuration/)
            pptx_files = list(mymastery_path.rglob(f"{lesson_code}*.pptx"))
            # Filter out temp files
            pptx_files = [f for f in pptx_files if not f.name.startswith("~")]
            if not pptx_files:
                return ""
            pptx_path = pptx_files[0]
        else:
            # Assume it's already a file path
            pptx_path = mymastery_path

        prs = Presentation(str(pptx_path))

        # Search multiple slides if slide_num not specified
        slide_nums = [4, 5, 6, 7] if slide_num is None else [slide_num]

        # Try each slide to find objectives
        for slide_num in slide_nums:
            if len(prs.slides) < slide_num:
                continue

            slide_index = slide_num - 1  # Convert to 0-indexed
            slide_text = []

            # Extract all text from the slide
            for shape in prs.slides[slide_index].shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                elif hasattr(shape, "text_frame") and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text.strip())

            # If we found text that looks like objectives, return it
            if slide_text:
                combined_text = "\n".join(slide_text)
                # Check if it contains objective-like content
                if any(
                    keyword in combined_text.lower()
                    for keyword in ["following this lesson", "students will", "objective", "be able to"]
                ):
                    return combined_text
                # If we have substantial text, use it
                if len(combined_text) > 50:
                    return combined_text

        return ""
    except Exception as e:
        print(f"Error extracting objectives from PowerPoint: {e}")
        return ""


def get_first_objective(objectives_text):
    """Extract only the first objective from text."""
    if not objectives_text:
        return ""

    # Split by line breaks and find first substantial line (>10 chars)
    lines = [line.strip() for line in objectives_text.split("\n") if line.strip()]
    for line in lines:
        if len(line) > 10:
            return line

    # Fallback: return first 200 characters
    return objectives_text[:200]
